"""
Convert files in uploads/ to markdown and route them to knowledge/ subfolders.

Supports: PDF, Word (.docx), PowerPoint (.pptx), CSV, Notion exports, markdown, txt.

Usage:
    python scripts/ingest.py
    python scripts/ingest.py --dry-run

Routing rules (filename keywords, case-insensitive):
    brand, voice, tone, guideline    -> knowledge/brand/
    icp, persona, audience           -> knowledge/icp/
    service, product, offering, sku  -> knowledge/services/
    case, study, customer, story     -> knowledge/content-library/case-studies/
    kpi, metric, dashboard, okr      -> knowledge/  (as kpis.md or appended)
    competitor, landscape, market    -> knowledge/markets/
    blog, post, article, newsletter  -> knowledge/content-library/
    everything else                  -> knowledge/content-library/inbox/

Dependencies (install on demand):
    pip install pypdf python-docx python-pptx
"""

from __future__ import annotations

import argparse
import csv
import re
import shutil
import sys
from datetime import datetime
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
UPLOADS = ROOT / "uploads"
KNOWLEDGE = ROOT / "knowledge"

ROUTES = [
    (("brand", "voice", "tone", "guideline"), "brand"),
    (("icp", "persona", "audience"), "icp"),
    (("service", "product", "offering", "sku"), "services"),
    (("case", "study", "customer", "story"), "content-library/case-studies"),
    (("competitor", "landscape", "market"), "markets"),
    (("blog", "post", "article", "newsletter"), "content-library"),
]


def slugify(name: str) -> str:
    s = re.sub(r"[^a-zA-Z0-9]+", "-", name.lower()).strip("-")
    return s or "untitled"


def route_for(filename: str) -> str:
    low = filename.lower()
    if any(k in low for k in ("kpi", "metric", "dashboard", "okr")):
        return "kpis"
    for keywords, folder in ROUTES:
        if any(k in low for k in keywords):
            return folder
    return "content-library/inbox"


def read_zip(path: Path) -> str:
    import zipfile
    import tempfile
    chunks = []
    with tempfile.TemporaryDirectory() as tmp:
        with zipfile.ZipFile(str(path), "r") as zf:
            zf.extractall(tmp)
        for child in sorted(Path(tmp).rglob("*")):
            if child.is_file() and child.suffix.lower() in (".md", ".txt", ".csv"):
                content = convert(child)
                if content:
                    chunks.append(f"## {child.name}\n\n{content}")
    return "\n\n---\n\n".join(chunks) if chunks else f"[no readable files found in {path.name}]"


def read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return f"[install pypdf to extract: {path.name}]"
    reader = PdfReader(str(path))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def read_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        return f"[install python-docx to extract: {path.name}]"
    doc = Document(str(path))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return f"[install python-pptx to extract: {path.name}]"
    prs = Presentation(str(path))
    chunks = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks.append(f"## Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                chunks.append(shape.text.strip())
    return "\n\n".join(chunks)


def read_csv(path: Path) -> str:
    rows = []
    with path.open(newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(" | ".join(row))
    return "\n".join(rows)


def read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


READERS = {
    ".pdf": read_pdf,
    ".docx": read_docx,
    ".pptx": read_pptx,
    ".csv": read_csv,
    ".md": read_text,
    ".txt": read_text,
    ".zip": read_zip,
}


def convert(path: Path) -> str | None:
    reader = READERS.get(path.suffix.lower())
    if not reader:
        return None
    try:
        return reader(path)
    except Exception as e:
        return f"[error extracting {path.name}: {e}]"


def write_markdown(content: str, source: Path, target_dir: Path) -> Path:
    target_dir.mkdir(parents=True, exist_ok=True)
    today = datetime.now().strftime("%d-%m-%Y")
    slug = slugify(source.stem)
    out = target_dir / f"{today}-{slug}.md"
    frontmatter = (
        f"---\n"
        f"source: {source.name}\n"
        f"ingested: {today}\n"
        f"---\n\n"
    )
    out.write_text(frontmatter + content, encoding="utf-8")
    return out


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--dry-run", action="store_true")
    args = ap.parse_args()

    if not UPLOADS.exists():
        print(f"No uploads/ folder. Create it and drop files in.")
        sys.exit(1)

    files = [p for p in UPLOADS.rglob("*") if p.is_file() and not p.name.startswith(".")]
    if not files:
        print("uploads/ is empty.")
        return

    print(f"Found {len(files)} file(s) in uploads/.\n")
    written = 0
    for path in files:
        route = route_for(path.name)
        if route == "kpis":
            target = KNOWLEDGE
            target_path = target / "kpis.md"
        else:
            target = KNOWLEDGE / route
            target_path = None

        content = convert(path)
        if content is None:
            print(f"  skip (unsupported): {path.name}")
            continue

        if args.dry_run:
            print(f"  would write: {path.name} -> {target.relative_to(ROOT)}/")
            continue

        if route == "kpis":
            target.mkdir(parents=True, exist_ok=True)
            today = datetime.now().strftime("%d-%m-%Y")
            existing = target_path.read_text(encoding="utf-8") if target_path.exists() else ""
            target_path.write_text(
                existing + f"\n\n## From {path.name} ({today})\n\n{content}\n",
                encoding="utf-8",
            )
            print(f"  appended: {path.name} -> knowledge/kpis.md")
        else:
            out = write_markdown(content, path, target)
            print(f"  wrote: {path.name} -> {out.relative_to(ROOT)}")
        written += 1

    if not args.dry_run:
        print(f"\nIngested {written} file(s). Originals remain in uploads/.")
        print("Review knowledge/ and move misrouted files manually if needed.")


if __name__ == "__main__":
    main()
